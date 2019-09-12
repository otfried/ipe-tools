#
# poweripe.py
#
# Translate an Ipe presentation to Powerpoint
#

import sys, os, re, datetime, argparse

import pptx
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR

# TODO: maybe use a better Latex parser
# from pylatexenc.latexwalker import LatexWalker, LatexEnvironmentNode

from PIL import Image
from lxml import etree

import ipe as Ipe
ipe = Ipe.ipe()

assert sys.hexversion >= 0x3060000

# --------------------------------------------------------------------

# these values are from 'official' Ipe presentation stylesheet
# TODO: compute based on actual stylesheet
sizeMap = { "huge": 2.8 * 17,
            "large": 2.8 * 12,
            "normal": 2.8 * 10,
            "small" : 2.4 * 10,
}

# Latex definitions are:
latex_sizes = {
  "\tiny", 5,
  "\footnotesize", 8,
  "\small", 9,
  "\normalsize", 10,
  "\large", 12,
  "\Large", 14, 
  "\LARGE", 17,
  "\huge", 20, 
  "\Huge", 25,
  }

def addBulletToParagraph(p):
  size = p.font.size or Pt(18)
  el = p._pPr
  el.set("marL", "%d" % size)
  el.set("indent", "-%d" % size)
  nsmap = el.nsmap
  el.insert(0, etree.Element("{" + nsmap['a'] + "}buChar", char="\u2022"))

# --------------------------------------------------------------------

class PowerIpe:
  def __init__(self, args):
    self.args = args

  def findColors(self):
    s = self.sheets.allNames(self.sheets, "color")
    self.colors = {}
    for k in s:
      sym = s[k]
      abs = self.sheets.find(self.sheets, "color", sym)
      h = "%02x%02x%02x" % (int(abs.r * 255), int(abs.g * 255), int(abs.b * 255))
      rgb = pptx.dml.color.RGBColor.from_string(h)
      self.colors[sym] = rgb
    
  def convertTime(self, s):
    return datetime.datetime(int(s[2:6]),
                             int(s[6:8]),
                             int(s[8:10]),
                             int(s[10:12]),
                             int(s[12:14]),
                             int(s[14:]))

  def addSvg(self, slide, pageNo, svgobjs):
    doc1 = ipe.Document()
    doc1.replaceSheets(doc1, self.sheets.clone(self.sheets))
    doc1.setProperties(doc1, self.iprops)
    p1 = doc1[1]
    p = self.doc[pageNo]
    for i in svgobjs:
      obj = p[i]
      p1.insert(p1, None, obj.clone(obj), None, "alpha")
    tmpname = "tmpipe%d.ipe" % pageNo
    svgname = "tmpipe%d.svg" % pageNo
    pngname = "tmpipe%d.png" % pageNo
    doc1.save(doc1, tmpname, "xml", None)
    os.system("iperender -svg -nocrop %s %s" % (tmpname, svgname))
    os.system("iperender -png -nocrop -transparent %s %s" % (tmpname, pngname))
    pic = slide.shapes.add_svg_picture(svgname, pngname, 0, 0)

  def convertText(self, slide, obj):
    s = obj.text(obj).strip()
    textStyle = obj.get(obj, "textstyle")
    textSize = obj.get(obj, "textsize")
    color = obj.get(obj, "stroke")
    italic = False
    bold = False
    if s.startswith(r"\bfseries"):
      s = s[9:].lstrip()
      bold = True
    if s.startswith(r"\itshape"):
      s = s[8:].lstrip()
      italic = True
    parSplitRe = "\n[ \t]*\n|\\\\par\\W*"
    if textStyle == "itemize" and s.startswith(r"\item"):
      s = s[5:].lstrip()
      parSplitRe += "|\\\\item"
      textStyle = "item"
    pars = re.split(parSplitRe, s)
    pars = [ s.replace("\n", " ") for s in pars ]
    if len(pars) == 0:
      return # huh?
    box = ipe.Rect()
    obj.addToBBox(obj, box, ipe.Matrix())
    pos = box.topLeft(box)
    size = sizeMap.get(textSize, sizeMap["normal"])
    #helper = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(self.x0 + pos.x), Pt(self.y0 - pos.y),
    # Pt(box.width(box)), Pt(box.height(box)))
    txBox = slide.shapes.add_textbox(Pt(self.x0 + pos.x), Pt(self.y0 - pos.y - 0.2 * size),
                                     Pt(box.width(box)), Pt(box.height(box)))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_top = 0
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i in range(len(pars)):
      s = pars[i].strip()
      s = re.sub(r"\\\\[ \t\n\r\f\v]*", "\v", s)
      if i == 0:
        p = tf.paragraphs[0]
      else:
        p = tf.add_paragraph()
      p.text = s
      p.font.size = Pt(size)
      p.line_spacing = Pt(1.2 * size)
      if textStyle == "center":
        p.alignment = PP_ALIGN.CENTER
      elif textStyle == "item":
        addBulletToParagraph(p)
      if color in self.colors:
        p.font.color.rgb = self.colors[color]
      p.font.bold = bold
      p.font.italic = italic
      if "\\textcolor{" in s:
        pos = 0
        r = p
        pat = re.compile("\v|\\\\textcolor{([a-zA-Z0-9]+)}{([^}]+)}")
        while True:
          m = pat.search(s, pos)
          if m is None:
            r.text = s[pos:]
            break
          r.text = s[pos:m.start()]
          if m.lastindex is None:
            p.add_line_break()
          else:
            r = p.add_run()
            r.text = m.group(2)
            runColor = m.group(1)
            if runColor in self.colors:
              r.font.color.rgb = self.colors[runColor]
          r = p.add_run()
          pos = m.end()

  # determine if we want to convert this text object
  def wantConvertText(self, obj):
    if self.args.no_text:
      return False
    mp = obj.get(obj, "minipage")
    if not mp:
      return False
    s = obj.text(obj)
    if self.args.latex:
      return True
    # allow some simple latex
    s = s.strip()
    if s.startswith("\\itshape"):
      s = s[8:]
    if s.startswith("\\bfseries"):
      s = s[9:]
    s = s.replace(r"\\", "")
    s = s.replace(r"\textcolor", "")
    s = s.replace(r"\item", "")
    s = s.replace(r"\par", "")
    return not ("\\" in s or "$" in s)

  def embedImage(self, slide, obj):
    info = obj.info(obj)
    s = str(obj) # make a unique name
    imgname = "tmpipe" + s[s.index("@"):]
    ext = ".jpg" if info.format == "jpg" else ".ppm"
    obj.savePixels(obj, imgname + ext)
    if ext == ".ppm":
      im = Image.open(imgname + ext)
      ext = ".png"
      im.save(imgname + ext)
    box = ipe.Rect()
    obj.addToBBox(obj, box, ipe.Matrix())
    pos = box.topLeft(box)
    pic = slide.shapes.add_picture(imgname + ext, Pt(self.x0 + pos.x), Pt(self.y0 - pos.y),
                                   Pt(box.width(box)), Pt(box.height(box)))
    
  def convertPage(self, pageNo, viewNo):
    sys.stderr.write("Converting page %d\n" % pageNo)
    p = self.doc[pageNo]
    t = p.titles(p)
    if t.title != "":
      slide = self.prs.slides.add_slide(self.title_slide_layout)
      title_placeholder = slide.shapes.title
      title_placeholder.text = t.title
    else:
      slide = self.prs.slides.add_slide(self.blank_slide_layout)
    # First determine which objects we'll convert to svg
    svgobjs = set()
    for i in range(1, len(p) + 1):
      if not p.visible(p, viewNo, i):
        continue
      obj = p[i]
      ty = obj.type(obj)
      if ty == "text" and self.wantConvertText(obj):
        pass
      elif ty == "image":
        self.embedImage(slide, obj)
      else:
        svgobjs.add(i)
    self.addSvg(slide, pageNo, svgobjs)
    # Now add remaining objects
    for i in range(1, len(p) + 1):
      if i in svgobjs or not p.visible(p, viewNo, i):
        continue
      obj = p[i]
      ty = obj.type(obj)
      if ty == "text":
        self.convertText(slide, obj)

  def convert(self):
    self.doc = ipe.Document(self.args.ipefile)
    self.sheets = self.doc.sheets(self.doc)
    self.iprops = self.doc.properties(self.doc)
    self.findColors()

    self.prs = pptx.Presentation()
    self.title_slide_layout = self.prs.slide_layouts[5]
    self.blank_slide_layout = self.prs.slide_layouts[6]

    props = self.prs.core_properties
    props.author = self.iprops.author
    props.last_modified_by = "" # shouldn't be python-pptx author
    props.title = self.iprops.title
    props.created = self.convertTime(self.iprops.created)
    props.modified = self.convertTime(self.iprops.modified)
    props.keywords = self.iprops.keywords
    props.subject = self.iprops.subject
    
    layout = self.sheets.find(self.sheets, "layout")
    self.x0 = layout.origin.x
    self.y0 = layout.origin.y + layout.papersize.y
  
    self.prs.slide_width = Pt(layout.papersize.x)
    self.prs.slide_height = Pt(layout.papersize.y)
  
    for pageNo in range(1, len(self.doc) + 1):
      p = self.doc[pageNo]
      viewNo = p.countViews(p)
      self.convertPage(pageNo, viewNo)

    pptname = self.args.output
    if pptname is None:
      pptname = self.args.ipefile[:-4] + '.pptx'
    self.prs.save(pptname)
  
# --------------------------------------------------------------------

if __name__ == "__main__":
  parser = argparse.ArgumentParser(description="Convert an Ipe presentation into pptx format")
  parser.add_argument('ipefile', help="File name for Ipe input")
  parser.add_argument('--output', help="File name for pptx output")
  parser.add_argument('--no-text', help="Do not convert text objects", action='store_true')
  parser.add_argument('--latex', help="Try to convert Latex text", action='store_true')

  powerIpe = PowerIpe(parser.parse_args())
  powerIpe.convert()

# --------------------------------------------------------------------
